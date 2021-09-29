"""
Classes is a source page for classes and objects during the program

The gui class is mandatory to all options - and it is used once.
Most of the essential methods in this program are driven from the class "standard analysis"

"""

# - internal imports:
from utils import *
from pptx.enum.text import PP_ALIGN

# Defines
define_data_analysis = "ניתוח נתונים"
define_split = "פיצול ענפים / מקצועות"
define_matrix = "טבלת הצלבות נתונים - מטריצה"
define_unite_files = "איחוד קבצים"


# class base
class GUIInput:
    # class attribute
    species = "root"

    # instance attributes
    def __init__(self, choice_area=None, choice_specific=None, input_file=None, second_input_file=None,
                 output_directory=None, filter_instructions_array=None, number_of_files=0):
        if filter_instructions_array is None:
            filter_instructions_array = [["ללא סינון", None]]
        self.choice_area = choice_area
        self.choice_specific = choice_specific
        self.input_file = input_file
        self.second_input_file = second_input_file
        self.output_directory = output_directory
        self.filter_instructions_array = filter_instructions_array
        self.number_of_files = number_of_files

    # instance methods
    def get_choice_tree(self):
        print("Activating Choice Tree decisions - GUI")
        # activating decision tree for GUI
        if self.choice_area == define_data_analysis:  # B
            self.data_analysis_window()
            pass
        if self.choice_area == define_split:  # C
            self.split_window()
            pass
        if self.choice_area == define_matrix:  # 9
            self.choice_specific = 9
            self.matrix_window()
            pass
        if self.choice_area == define_unite_files:  # 10
            self.choice_specific = 10
            self.pop_up_num_files_window()
            self.unite_files_window()
            pass

    def welcome_window(self):
        ww = base_frame("מסך הגדרות למשתמש - תוכנת עיבוד נתונים")  # ww = welcome window
        button(ww, 15, 0, "Next", partial(move_to_window, self, ww, "get_choice_tree"))
        button(ww, 15, 1, "Exit", partial(move_to_window, self, ww, "exit"))
        label(ww, ",ברוך הבא לתוכנה האולטימטיבית של מחוז דרום בשירות התעסוקה\n .אנא בחר את האופציה המתאימה לך מטה", 10,
              1, 0, 0, 10, "#2B327A", "orange", E, 30)
        # prep result to present user
        res = Label(ww, fg="white", bg="#2B327A")
        res.grid(columnspan=10, row=10, column=0, sticky=S + E + N, padx=25, pady=5)
        res.config(text="\n\n\n\n", justify=RIGHT)
        # Radiobutton
        choice = IntVar()
        radio_button(ww, 2, 9, define_data_analysis, choice, 1,
                     partial(decisions_area_assignment, self, define_data_analysis, res))
        radio_button(ww, 3, 9, define_split, choice, 2,
                     partial(decisions_area_assignment, self, define_split, res))
        radio_button(ww, 4, 9, define_matrix, choice, 3,
                     partial(decisions_area_assignment, self, define_matrix, res))
        radio_button(ww, 6, 9, define_unite_files, choice, 5,
                     partial(decisions_area_assignment, self, define_unite_files, res))
        # add image IES
        img = PhotoImage(file=r"src_files/icon_SD.png").subsample(3, 3)
        Label(ww, image=img, bg="#2B327A").grid(rowspan=2, row=12, column=0, columnspan=10, padx=5, pady=35, sticky=N)
        ww.protocol("WM_DELETE_WINDOW", exit)
        ww.mainloop()  # run the window endlessly until user response


    def data_analysis_window(self):
        dw = base_frame("עיבוד נתוני שירות התעסוקה - אפשרויות")  # dw = data window
        button(dw, 15, 8, "Back", partial(move_to_window, self, dw, "welcome_window"))
        label(dw, "\n :אנא לחץ על האופציה המתאימה לך מטה", 10, 1, 0, 0, 10, "#2B327A", "orange", E, 30)
        # buttons
        button(dw, 2, 0, "(מומלץ) ניתוח נתונים סטנדרטי לפי פורמט השירות",
               partial(data_analysis_button_reactor, dw, self, 1), "black", "light blue", 70, 5, 5, 9)
        button(dw, 3, 0, "ניתוח נתונים סטנדרטי לפי פורמט אוטומטי", partial(data_analysis_button_reactor, dw, self, 8),
               "black", "light blue", 70, 5, 5, 9)
        button(dw, 4, 0, "ניתוח נתונים על פי רשימת ת''ז לפי פורמט השירות",
               partial(data_analysis_button_reactor, dw, self, 4), "black", "light blue", 70, 5, 5, 9)
        button(dw, 5, 0, "ניתוח נתונים סטנדרטי לכל לשכות מחוז דרום", partial(data_analysis_button_reactor, dw, self, 2),
               "black", "light blue", 70, 5, 5, 9)
        button(dw, 6, 0, "ניתוח נתונים סטנדרטי לכל מחוזות השירות", partial(data_analysis_button_reactor, dw, self, 3),
               "black", "light blue", 70, 5, 5, 9)
        # add image IES
        img = PhotoImage(file=r"src_files/icon_SD.png").subsample(3, 3)
        Label(dw, image=img, bg="#2B327A").grid(rowspan=2, row=12, column=0, columnspan=10, padx=5, pady=60, sticky=N)
        dw.protocol("WM_DELETE_WINDOW", exit) # if user exit- program ends
        dw.mainloop()  # run the window endlessly until user response

    def data_analysis_input_window(self):
        dw = base_frame("ניתוח נתונים - הזנת ערכים וקבצים לתוכנה", 11, 6)  # dw = data window
        tmp = ["0", "ניתוח נתונים סטנדרטי לפי פורמט שירות התעסוקה",
               "ניתוח נתונים מחוזי לפי לשכות מחוז דרום",
               "ניתוח נתונים ארצי לפי מחוזות הארץ",
               "ניתוח נתונים סטנדרטי לפי הכנסת רשימת תעודות זהות",
               "5", "6", "7", "ניתוח נתונים סטנדרטי לפי פורמט אוטומטי"]
        text = "זהו חלון ההגדרות עבור ניתוח נתונים לפי דוח ממחולל הדוחות של שירות התעסוקה\nעל מנת להשתמש " \
               "בפונקציית זו יש תחילה להוציא דוח ממחולל הדוחות של השירות\n"
        specified_text = text + "\n" + "התוכנית שנבחרה: " + tmp[self.choice_specific]
        label(dw, specified_text, 6, 1, 0, 5, 9, "#35B7E8", "Black", S + W + E + N, 5, 0)
        # upload a file - user
        label(dw, "אנא בחרו קובץ נתונים (ייצוא ממחולל הדוחות)", 4, 2, 2, 1, 10, '#2B327A', "#E98724", NE, 15, 1)
        button(dw, 3, 2, "לחץ כאן לבחירת קובץ", partial(choose_file, self), "white", "black", None, 15, 1, 4)
        # choose a folder destination
        label(dw, "אנא בחרו תקיית יעד לתוצרי המערכת", 4, 4, 2, 1, 10, '#2B327A', "#E98724", NE, 15, 1)
        button(dw, 5, 2, "לחץ כאן לבחירת תקייה", partial(choose_output_path_folder, self), "white", "black", None, 15,
               1, 4)
        if self.choice_specific == 1 or self.choice_specific == 8:
            # choose filter groups
            button(dw, 9, 2, "לחץ כאן לבחירת קבוצות מיקוד לסינון (אופציונלי)",
                   partial(move_to_window, self, None, "filter_group_user_input_window"),
                   "black", "#35B7E8", None, 15, 80, 4)
        if self.choice_specific == 4:
            # choose ID file
            button(dw, 8, 2, "לחץ כאן לבחירת קובץ עם רשימת תעודות הזהות שברצונך לסנן",
                   partial(choose_file, self, True),
                   "black", "orange", None, 15, 30, 4)
            # choose filter groups
            button(dw, 9, 2, "לחץ כאן לבחירת קבוצות מיקוד לסינון (אופציונלי)",
                   partial(move_to_window, self, None, "filter_group_user_input_window"),
                   "black", "#35B7E8", None, 15, 30, 4)
        button(dw, 10, 0, "Start", partial(move_to_window, self, dw, "check n close"))
        button(dw, 10, 1, "Back", partial(move_to_window, self, dw, "data_analysis_window"))
        dw.protocol("WM_DELETE_WINDOW", exit) # if user exit- program ends
        dw.mainloop()  # run the window endlessly until user response

    def split_window(self):
        sw = base_frame("פונציית הפיצול של איציק - ענפים ומקצועות", 11, 6)  # sw = split window
        text = "זהו חלון ההגדרות עבור פיצול עמודת מקצועות או ענפים בדוח מחולל הדוחות של שירות התעסוקה\nעל מנת להשתמש " \
               "בפונקציית הפיצול של איציק יש תחילה להוציא דוח ממחולל הדוחות של השירות\nולוודא כי העמודה אותה נרצה " \
               "לפצל קיימת בדוח "
        label(sw, text, 6, 1, 0, 5, 9, "#35B7E8", "Black", S + W + E + N, 5, 0)
        # upload a file - user
        label(sw, "אנא בחרו קובץ נתונים (ייצוא ממחולל הדוחות)", 4, 2, 2, 1, 10, '#2B327A', "#E98724", NE, 15, 1)
        button(sw, 3, 2, "לחץ כאן לבחירת קובץ", partial(choose_file, self), "white", "black", None, 15, 1, 4)
        # choose a folder destination
        label(sw, "אנא בחרו תקיית יעד לתוצרי המערכת", 4, 4, 2, 1, 10, '#2B327A', "#E98724", NE, 15, 1)
        button(sw, 5, 2, "לחץ כאן לבחירת תקייה", partial(choose_output_path_folder, self), "white", "black", None, 15,
               1, 4)
        # mark options to split
        jobs = IntVar()
        fields = IntVar()
        Label(sw, text=":בחירת עמודות לפיצול", fg="#FFFFFF", font=("Arial", 11, "bold"), justify=RIGHT, bg="#2B327A"). \
            grid(row=6, column=2, columnspan=4, padx=15, pady=15, sticky=SE)
        check_button(sw, 7, 5, jobs, 1, "פיצול עמודת מקצועות")
        check_button(sw, 8, 5, fields, 0, "פיצול עמודת ענפים")
        # choose filter groups
        button(sw, 9, 2, "לחץ כאן לבחירת קבוצות מיקוד לסינון (אופציונלי)",
               partial(move_to_window, self, None, "filter_group_user_input_window"),
               "black", "#35B7E8", None, 15, 80, 4)
        button(sw, 10, 0, "Start", partial(check_box_for_split, self, jobs, fields, sw))
        button(sw, 10, 1, "Back", partial(move_to_window, self, sw, "welcome_window"))
        sw.protocol("WM_DELETE_WINDOW", exit) # if user exit- program ends
        sw.mainloop()  # run the window endlessly until user response

    def matrix_window(self):
        mw = base_frame("הצלבת נתונים - יצירת מטריצת נתונים", 12, 6)  # mw = matrix window
        button(mw, 11, 0, "Next", partial(move_to_window, self, mw, "check n close"))
        button(mw, 11, 1, "Back", partial(move_to_window, self, mw, "welcome_window"))
        text = ".זהו חלון ההגדרות עבור הצלבת נתוני דורשי עבודה מדוח של מחולל הדוחות של שירות התעסוקה\nעל מנת להשתמש " \
               "בפונקצייה זו יש תחילה להוציא דוח ממחולל הדוחות של השירות לוודא כי\n.העמודות אותן נרצה " \
               "להצליב קיימות בדוח\n.אנא שימו לב כי ערכים ייחודים שאותם לא ניתן לסכום לא יכנסו לטבלת הנתונים"
        label(mw, text, 6, 1, 0, 5, 9, "#35B7E8", "Black", S + W + E + N, 5, 0)
        # upload a file - user
        label(mw, "אנא בחרו קובץ נתונים (ייצוא ממחולל הדוחות)", 4, 2, 2, 1, 10, '#2B327A', "#E98724", NE, 15, 1)
        button(mw, 3, 2, "לחץ כאן לבחירת קובץ", partial(choose_file, self), "white", "black", None, 15, 1, 4)
        # choose a folder destination
        label(mw, "אנא בחרו תקיית יעד לתוצרי המערכת", 4, 4, 2, 1, 10, '#2B327A', "#E98724", NE, 15, 1)
        button(mw, 5, 2, "לחץ כאן לבחירת תקייה", partial(choose_output_path_folder, self), "white", "black", None, 15,
               1, 4)
        # choose filter groups
        button(mw, 9, 2, "לחץ כאן לבחירת קבוצות מיקוד לסינון (אופציונלי)",
               partial(move_to_window, self, None, "filter_group_user_input_window"),
               "black", "#35B7E8", None, 15, 140, 4)
        mw.protocol("WM_DELETE_WINDOW", exit) # if user exit- program ends
        mw.mainloop()  # run the window endlessly until user response

    def pop_up_num_files_window(self):
        self.number_of_files = 1
        w = tk.Tk()
        w.minsize(320, 80)
        try:
            w.iconphoto(True, tk.PhotoImage(file='src_files/icon.png'))
        except:
            print("error in uploading tk.PhotoImage in window")
        w.title("E2P")
        # opening statement
        tk.Label(w, text="אנא בחרו את כמות הקבצים שתרצו לאחד", justify=RIGHT).grid(row=0, column=3, columnspan=3,
                                                                                   sticky=SE)
        listbox = tk.Listbox(w, height=5, width=12, bg="#FFFFFF", activestyle='dotbox', fg="#2B327A", justify=LEFT)
        listbox.insert(1, "1")
        listbox.insert(2, "2")
        listbox.insert(3, "3")
        listbox.insert(4, "4")
        listbox.insert(5, "5")
        listbox.grid(columnspan=3, row=1, column=3, sticky=NE, padx=5, pady=5)
        button(w, 2, 1, "Next", partial(retrieve, w, listbox, self), "black", "#a9a9a9")
        w.protocol("WM_DELETE_WINDOW", exit) # if user exit- program ends
        w.mainloop()

    def unite_files_window(self):
        print("Chose to combine " + str(self.number_of_files) + " files. uploading window for uploading.")
        self.input_file = []
        uw = base_frame("איחוד קבצים - יצירת קובץ משותף", 12, 6)  # uw = unite window
        button(uw, 11, 0, "Start", partial(move_to_window, self, uw, "check n close"))
        button(uw, 11, 1, "Back", partial(move_to_window, self, uw, "welcome_window"))
        text = ".זהו חלון ההגדרות עבור איחוד קבצים שונים של נתוני דורשי עבודה מדוח של מחולל הדוחות של שירות התעסוקה" + "\n" \
                                                                                                                       "על מנת להשתמש בפונקצייה זו יש תחילה להוציא דוחות מפוצלים מהמחולל - לדוגמא דוח עם ערכים זהים עבור כל" + "\n" \
                                                                                                                                                                                                                               ".מחוז בנפרד בכדי לאחדם לקובץ שיצור נתונים של כלל הארץ. מטרת הפונקצייה היא שימוש בנתונים רבים בבת אחת"
        label(uw, text, 6, 1, 0, 5, 9, "#35B7E8", "Black", S + W + E + N, 5, 0)
        # upload a file - user
        label(uw, "אנא בחרו קבצי נתונים (ייצוא ממחולל הדוחות)", 4, 2, 2, 1, 10, '#2B327A', "#E98724", NE, 15, 10)
        for i in range(int(self.number_of_files)):
            txt = "לחץ כאן לבחירת קובץ מספר " + str(i + 1)
            button(uw, 3 + i, 2, txt, partial(choose_multiple_files, self), "white", "black", None, 15, 1, 4)
        # choose a folder destination
        label(uw, "אנא בחרו תקיית יעד לתוצרי המערכת", 4, 8, 2, 1, 10, '#2B327A', "#E98724", NE, 15, 1)
        button(uw, 9, 2, "לחץ כאן לבחירת תקייה", partial(choose_output_path_folder, self), "black", "gray", None, 15, 1,
               4)
        # choose filter groups
        button(uw, 10, 2, "לחץ כאן לבחירת קבוצות מיקוד לסינון (אופציונלי)",
               partial(move_to_window, self, None, "filter_group_user_input_window"),
               "black", "#35B7E8", None, 15, 50, 4)
        uw.protocol("WM_DELETE_WINDOW", exit) # if user exit- program ends
        uw.mainloop()  # run the window endlessly until user response

    def filter_group_user_input_window(self):
        print("Uploading filter group choosing window... GUI from user.")
        # window
        fw = base_frame("בחירת קבוצות מיקוד - לסינון קובץ הנתונים", top_lvl_flag=True)  # fw = filter window
        fw.configure(bg="#35B7E8")
        text = ".זהו חלון המיועד להזנת קבוצות מיקוד נוספות לסינון מתוך המאגר הסטטיסטי שהכנסתם למערכת בחלון הקודם" + "\n" + \
               "אנא שימו לב לדוגמא מטה והזינו את ערכי הסינון בהתאם להנחיות. שימו לב לא להכניס סימני פיסוק או ערכים" + "\n" + \
               "שלא קיימים בקובץ הנתונים שהזנתם בחלון הקודם. שימו לב! כדי להזין ערכי סינון מרובים יש להפרידם באמצעות" + "\n" + \
               ".הסימון ',' (פסיק) בלבד. קיימות לבחירתכם 3 אפשרויות בסיס לבחירה לבחירות השכיחות ביותר"
        label(fw, text, 10, 1, 0, 6, 9, "#35B7E8", "Black", S + W + E + N, 5, 0)
        label(fw, ":בחירת קבוצות סינון מוגדרות מראש", 10, 2, 0, 0, 11, "white", "#E98724", S + E, 15, 1)
        # choose default filters:
        county_default_filter = IntVar()  # all country
        district_default_filter = IntVar()  # all district
        office_choice_filter = IntVar()  # pick an office
        check_button(fw, 3, 9, county_default_filter, 1, "כלל הארץ (ללא סינון כלל)", 1, "#35B7E8")
        check_button(fw, 4, 9, district_default_filter, 1, "מחוז דרום", 1, "#35B7E8")
        check_button(fw, 5, 9, office_choice_filter, 0, ":לשכה - אנא הזן שם (לדוגמא 'אופקים')", 3, "#35B7E8")
        office_name_user_input = Entry(fw, fg="#2B327A", width=20, justify=RIGHT)
        office_name_user_input.grid(row=5, column=0, columnspan=6, sticky=E, padx=2)
        # prepare more values - through pickle
        headers = ['לשכה', 'מחוז', 'מגדר', 'גיל', 'מצב משפחתי',
                   'ילדים עד גיל 18', 'סוג תביעה נוכחי', 'רמת השכלה', 'סיבת רישום']
        # choose more filters
        label(fw, ":בחירת קבוצות סינון נוספות", 10, 6, 0, 0, 11, "white", "#E98724", S + E, 15, 1)
        # 1
        label(fw, "קבוצה #1: הזן שם קבוצה בתיבה מטה", 3, 7, 7, 0, 9, "#35B7E8", "black", S + E, 15, 1)
        group1_name_user_input = Entry(fw, fg="#2B327A", width=30, justify=RIGHT)
        group1_name_user_input.grid(row=8, column=0, columnspan=10, sticky=E, padx=20, pady=5)
        head11 = list_box(fw, 9, 6, 4, "header", headers)
        l11 = list_box(fw, 9, 0, 6, "options", ["pick a field first"])
        head12 = list_box(fw, 10, 6, 4, "header", headers)
        l12 = list_box(fw, 10, 0, 6, "options", ["pick a field first"])
        head11.bind("<<ComboboxSelected>>", partial(update_filter_values, head11, l11))
        head12.bind("<<ComboboxSelected>>", partial(update_filter_values, head12, l12))
        # 2
        label(fw, "קבוצה #2: הזן שם קבוצה בתיבה מטה", 3, 11, 7, 0, 9, "#35B7E8", "black", S + E, 15, 1)
        group2_name_user_input = Entry(fw, fg="#2B327A", width=30, justify=RIGHT)
        group2_name_user_input.grid(row=12, column=0, columnspan=10, sticky=E, padx=20, pady=5)
        head21 = list_box(fw, 13, 6, 4, "header", headers)
        l21 = list_box(fw, 13, 0, 6, "options", ["pick a field first"])
        head22 = list_box(fw, 14, 6, 4, "header", headers)
        l22 = list_box(fw, 14, 0, 6, "options", ["pick a field first"])
        head21.bind("<<ComboboxSelected>>", partial(update_filter_values, head21, l21))
        head22.bind("<<ComboboxSelected>>", partial(update_filter_values, head22, l22))
        # 3
        label(fw, "קבוצה #3: הזן שם קבוצה בתיבה מטה", 3, 15, 7, 0, 9, "#35B7E8", "black", S + E, 15, 1)
        group3_name_user_input = Entry(fw, fg="#2B327A", width=30, justify=RIGHT)
        group3_name_user_input.grid(row=16, column=0, columnspan=10, sticky=E, padx=20, pady=1)
        group3_filter_header = Entry(fw, fg="BLACK", width=20, justify=RIGHT)
        group3_filter_header.insert(0, "ערך חופשי - כותרת")
        group3_filter_header.grid(row=17, column=6, columnspan=4, sticky=E, padx=20, pady=1)
        group3_filter_values = Entry(fw, fg="BLACK", width=35, justify=RIGHT)
        group3_filter_values.insert(0, "(ערך חופשי - ערכים (מופרדים בפסיק")
        group3_filter_values.grid(row=17, column=0, columnspan=6, sticky=E, padx=5, pady=1)
        """ set the values """
        button(fw, 20, 0, "Save", partial(set_filter_instructions_array, self, fw, county_default_filter,
                                          district_default_filter,
                                          office_choice_filter, office_name_user_input,
                                          group1_name_user_input, head11, l11, head12, l12,
                                          group2_name_user_input, head21, l21, head22, l22,
                                          group3_name_user_input, group3_filter_header,
                                          group3_filter_values))
        button(fw, 20, 1, "Cancel", fw.destroy)
        fw.protocol("WM_DELETE_WINDOW", exit) # if user exit- program ends
        fw.mainloop()  # run the window endlessly until user response

    def print_data_to_user(self):
        print("User have chose option number " + str(self.choice_specific) + " in area " + str(self.choice_area))
        print("input excel: " + str(self.input_file))
        print("output path folder: " + str(self.output_directory))
        print("filter instructions array contains " + str(len(self.filter_instructions_array)) + " study groups.")


class StandardAnalysis:
    # class attribute
    species = define_data_analysis

    # instance attributes
    def __init__(self, sheet_pd=None, filter_instructions_array=None, output_directory=None, query_table=None,
                 jobs_dic=None, fields_jobs_dic=None):
        self.sheet_pd = sheet_pd
        self.filter_instructions_array = filter_instructions_array
        self.output_directory = output_directory
        self.query_table_numbers = query_table
        self.query_table_percents = query_table
        self.jobs_dic = jobs_dic
        self.fields_jobs_dic = fields_jobs_dic

    def get_dictionary(self, param):
        """
            This function tops the Itzik function and after splitting - summing to a dictionary
            :param param: either jobs or fields (hebrew string)
            :return: a python dataframe that came from  dictionary that sums the param through all of the dataframe
        """
        if param in self.sheet_pd:
            print("")  # create separation
            sheet = get_splitted_sheet(self.sheet_pd, param)
            # create a comparison of job distributions
            heads = []
            l = 0
            output_dic = {}
            for col in self.filter_instructions_array:
                heads.append(col[0])  # create array heads for focus groups
                temp_filtered_df = sheet_pd_filter(sheet, col[1])  # filter dataframe
                # keep only columns with jobs
                pool_jobs = temp_filtered_df[
                    [str(param) + " " + str(1), str(param) + " " + str(2), str(param) + " " + str(3),
                     str(param) + " " + str(4)]]
                # create a dictionary {key = str of job : value = arr of int length of focus groups }
                output_dic = update_dic_values(output_dic, pool_jobs, l, len(self.filter_instructions_array))
                l = l + 1  # add to iterator
            # return a data frame
            if param == "מקצועות רלוונטיים":
                self.jobs_dic = pd.DataFrame.from_dict(output_dic, orient='index', columns=heads)
            elif param == "ענפי מקצועות רלוונטיים":
                self.fields_jobs_dic = pd.DataFrame.from_dict(output_dic, orient='index', columns=heads)
        else:
            print("No " + str(param) + " in the input excel sheet. Not activating the Itzik function.")
            return None

    def set_query_tables(self):
        print("Counting query data, setting query tables.")
        # prepare to count - add rows - IES format!
        rows_names = [
            None, "מגדר",
            "נקבה", "זכר",
            None, "סוג תביעה",
            "אבטלה", "הבטחת הכנסה", "אינו תובע",
            None, "סיבת רישום",
            "פיטורין", "חל''ת", "אינו עובד ומחפש עבודה", "התפטרות",
            None, "גיל",
            "15-24", "25-29", "30-34", "35-39", "40-44", "45-49", "50-54", "55-59", "60-64", "65-69", "70 +",
            None, "מספר ילדים עד גיל 18",
            "0", "1-2", "3-5", "6-8", "יותר מ-8",
            None, "מצב משפחתי",
            "רווק/ה", "נשוי/ה", "גרוש/ה",
            None, "מצב השכלה",
            "ללא השכלה", "תעודת בגרות", "יסודי/תיכון", "תואר ראשון", "תואר שני", "תואר שלישי", "תעודת מקצוע",
            None, None,
            "סכום כלל דורשי עבודה"
        ]
        output_num = pd.DataFrame({"אלמנט השוואה": rows_names}, columns=["אלמנט השוואה"])
        output_per = pd.DataFrame({"אלמנט השוואה": rows_names}, columns=["אלמנט השוואה"])
        # start counting - go over database {no. of columns} times
        for col in self.filter_instructions_array:
            name_col = col[0]
            temp_filtered_df = sheet_pd_filter(self.sheet_pd, col[1])  # filter dataframe
            # add output to
            temp_output = query_counter_helper(temp_filtered_df, name_col)  # counter helper
            output_num[name_col] = temp_output[0]
            output_per[name_col] = temp_output[1]
        self.query_table_numbers = output_num
        self.query_table_percents = output_per

    def create_graphs(self):
        print("Creating Graphs for the powerpoint presentation- output will be in a folder named 'Graphs'")
        query_sum_arr_for_graphs = []  # result tables
        # create a new directory
        try:
            os.mkdir(self.output_directory + '/Graphs')
        except FileExistsError:
            pass
        print("Directory 'Graphs' created")

        # prep Color Palette for the focus groups- taken from https://www.color-hex.com/color-palette/113194
        no_of_groups = len(self.filter_instructions_array)
        color_palette = ["#00aedb", "#00b159", "#d11141", "#f37735", "#ffc425", "#b0afac", "#902ef6", "#f62eee",
                         "#12a197", "#c1ebe0", "#b888be", "#c1e1be", "#99afd7", "#c18298", "#de9c66", "#96b86f",
                         "#96c1c7", "#a07db0", "#000000", "#808080", "#ffffff", "#df0f4d", "#d3d3d3", "#997761"]
        color_palette = color_palette[0:no_of_groups]

        # Bold font
        plt.rcParams["font.weight"] = "bold"
        plt.rcParams["axes.labelweight"] = "bold"

        """gender"""
        if "מגדר" in self.sheet_pd:
            filtered_sheet = self.query_table_percents.copy()  # go over the df we made and brake it to tables
            filtered_sheet = filtered_sheet[
                (filtered_sheet['אלמנט השוואה'] == 'נקבה') | (filtered_sheet['אלמנט השוואה'] == 'זכר')].T
            filtered_sheet = filtered_sheet.rename(columns=filtered_sheet.iloc[0])  # move first row to header
            filtered_sheet = filtered_sheet.iloc[1:]  # remove the first row
            query_sum_arr_for_graphs.append(filtered_sheet)  # for result
            filtered_sheet = get_hebrew_translation(filtered_sheet)  # Hebrew translation
            # graph
            graph_gender = filtered_sheet.plot.barh(stacked=True, color={"הבקנ": "#99FF99", "רכז": "#7b9bcc"})
            # attach values
            i = 0
            for bar in filtered_sheet.itertuples():
                f_pos = int(bar.הבקנ * 100)
                m_pos = int(bar.רכז * 100)
                if f_pos + m_pos != 100: m_pos = 100 - f_pos  # normalize
                graph_gender.text(f_pos / 200, i, str(f_pos) + '%', fontweight='bold')  # add values to graph
                graph_gender.text(f_pos / 100 + (m_pos / 200), i, str(m_pos) + '%',
                                  fontweight='bold')  # add values to graph
                i = i + 1
            graph_gender.set_xticks([])
            plt.legend(bbox_to_anchor=(1.25, 1), loc='upper right')
            plt.title('התפלגות מגדרית'[::-1], fontweight='bold')
            plt.xlim(0, 1)  # set x axis limit 1 - adjust frame
            plt.savefig(self.output_directory + '/Graphs/' + 'גרף_מגדר' + '.png',
                        bbox_inches='tight')  # save to folder as .png
            plt.clf()

        """Reason of registration"""
        if "סיבת רישום" in self.sheet_pd:
            filtered_sheet = self.query_table_percents.copy()  # go over the df we made and brake it to tables
            filtered_sheet = filtered_sheet[
                (filtered_sheet['אלמנט השוואה'] == 'פיטורין') | (
                        filtered_sheet['אלמנט השוואה'] == "חל''ת") | (
                        filtered_sheet['אלמנט השוואה'] == 'אינו עובד ומחפש עבודה') | (
                        filtered_sheet['אלמנט השוואה'] == 'התפטרות')]
            filtered_sheet.set_index('אלמנט השוואה', inplace=True)
            try:
                filtered_sheet.reindex(["חל''ת", "אינו עובד ומחפש עבודה", "התפטרות", "פיטורין"])  # re-order headlines
            except:
                pass
            query_sum_arr_for_graphs.append(filtered_sheet)  # for result
            filtered_sheet = get_hebrew_translation(filtered_sheet)  # Hebrew translation
            # graph
            graph_sue_type = filtered_sheet.plot.barh(figsize=(10, 6), color=color_palette)
            plt.title('סיבת רישום לשירות התעסוקה'[::-1], fontweight='bold')
            plt.ylabel("")
            max_value = max(filtered_sheet.max())  # get the max value in this dataframe
            plt.xlim(0, max_value + 0.05)  # set y axis limit 1 - adjust frame
            plt.legend(bbox_to_anchor=(1.2, 1), loc='upper right')  # legend outside the graph
            # set size of font
            if len(self.filter_instructions_array) > 3:
                size = 7
            else:
                size = 9
            # attach values
            for p in graph_sue_type.patches:
                if p.get_width() * 100 < 1:
                    h = "{:.1%}".format(p.get_width())
                else:
                    h = "{:.0%}".format(p.get_width())
                graph_sue_type.annotate(str(h), (p.get_width() + 0.005, p.get_y() + 0.06), size=size, rotation=0)
            plt.tick_params(axis='x', left=False, right=False, labelleft=False, labelbottom=False, bottom=False)
            plt.savefig(self.output_directory + '/Graphs/' + 'גרף_סיבת_רישום' + '.png',
                        bbox_inches='tight')  # save to folder as .png
            plt.clf()

        """Type of sue"""
        if "סוג תביעה נוכחי" in self.sheet_pd:
            filtered_sheet = self.query_table_percents.copy()  # go over the df we made and brake it to tables
            filtered_sheet = filtered_sheet[
                (filtered_sheet['אלמנט השוואה'] == 'אבטלה') | (
                        filtered_sheet['אלמנט השוואה'] == 'הבטחת הכנסה') | (
                        filtered_sheet['אלמנט השוואה'] == 'משותף') | (
                        filtered_sheet['אלמנט השוואה'] == 'אינו תובע')]
            filtered_sheet.set_index('אלמנט השוואה', inplace=True)
            colors = ['#6495ED', '#F4A460', '#90EE90', '#F08080', '#FFFF00', '#8FBC8F', '#66CDAA']
            query_sum_arr_for_graphs.append(filtered_sheet)  # for result
            i = 0  # create a different pie chart
            for group in filtered_sheet.columns:
                tmp = filtered_sheet.filter(items=[group])
                # graph creation
                y = tmp[group].to_numpy()  # get values
                mylabels = list(tmp.index.values)  # crate labels
                mylabels_bckword = []
                for word in mylabels: mylabels_bckword.append(word[::-1])  # Hebrew translation
                try:
                    _, labels, _ = plt.pie(y, labels=mylabels_bckword, colors=colors, shadow=True, autopct='%1.0f%%',
                                           textprops={'fontsize': 18}, normalize=False)  # create the pie
                except ValueError:
                    _, labels, _ = plt.pie(y, labels=mylabels_bckword, colors=colors, shadow=True, autopct='%1.0f%%',
                                           textprops={'fontsize': 18}, normalize=True)  # create the pie
                for lab in labels:
                    lab.set_fontsize(16)
                plt.title(("סוג תביעה: " + group)[::-1], fontweight='bold')
                plt.savefig(self.output_directory + '/Graphs/' + 'גרף_סוג_תביעה_' + str(group) + '.png',
                            bbox_inches='tight')  # save to folder as .png
                i += 1
                plt.clf()

        """Age"""
        if "גיל" in self.sheet_pd:
            filtered_sheet = self.query_table_percents.copy()  # go over the df we made and brake it to tables
            filtered_sheet = filtered_sheet[
                (filtered_sheet['אלמנט השוואה'] == '15-24') | (
                        filtered_sheet['אלמנט השוואה'] == '25-29') | (
                        filtered_sheet['אלמנט השוואה'] == '30-34') | (
                        filtered_sheet['אלמנט השוואה'] == '35-39') | (
                        filtered_sheet['אלמנט השוואה'] == '40-44') | (
                        filtered_sheet['אלמנט השוואה'] == '45-49') | (
                        filtered_sheet['אלמנט השוואה'] == '50-54') | (
                        filtered_sheet['אלמנט השוואה'] == '55-59') | (
                        filtered_sheet['אלמנט השוואה'] == '60-64') | (
                        filtered_sheet['אלמנט השוואה'] == '65-69') | (
                        filtered_sheet['אלמנט השוואה'] == '70 +')]
            filtered_sheet.set_index('אלמנט השוואה', inplace=True)
            query_sum_arr_for_graphs.append(filtered_sheet)  # for result
            # build graph
            x = ["18-24", "25-29", "30-34", "35-39", "40-44", "45-49", "50-54", "55-59", "60-64", "65-69", "70 +"]
            plt.figure(figsize=(10, 6))
            pos = 0
            for col in filtered_sheet.columns:
                tmp = filtered_sheet.filter(items=[col])
                # graph creation
                y = tmp[col].to_numpy()  # get values
                plt.plot(x, y, label=col[::-1], linewidth=4, color=color_palette[pos])  # add line to graph
                pos += 1
            plt.grid(linestyle='--', linewidth=0.8, which='both')  # add grid
            plt.legend(loc='upper right')  # legend outside the graph
            plt.title("התפלגות גילאי דורשי העבודה"[::-1], fontweight='bold')
            plt.gca().yaxis.set_major_formatter(ticker.PercentFormatter(1))  # manipulate to percents
            plt.savefig(self.output_directory + '/Graphs/'
                                                '' + 'גרף_גילאים' + '.png',
                        bbox_inches='tight')  # save to folder as .png
            plt.clf()

        """Children"""
        if "ילדים עד גיל 18" in self.sheet_pd:
            filtered_sheet = self.query_table_percents.copy()  # go over the df we made and brake it to tables
            filtered_sheet = filtered_sheet[
                (filtered_sheet['אלמנט השוואה'] == '0') | (
                        filtered_sheet['אלמנט השוואה'] == '1-2') | (
                        filtered_sheet['אלמנט השוואה'] == '3-5') | (
                        filtered_sheet['אלמנט השוואה'] == '6-8') | (
                        filtered_sheet['אלמנט השוואה'] == 'יותר מ-8')]
            filtered_sheet.set_index('אלמנט השוואה', inplace=True)
            query_sum_arr_for_graphs.append(filtered_sheet)  # for result
            filtered_sheet = get_hebrew_translation(filtered_sheet)  # Hebrew translation
            # graph
            graph_sue_type = filtered_sheet.plot.bar(rot=0, figsize=(11, 6), color=color_palette)
            plt.title('ילדים מתחת לגיל 81'[::-1], fontweight='bold')
            plt.grid(axis='y', linewidth=0.5, zorder=3)
            graph_sue_type.set_axisbelow(True)
            graph_sue_type.yaxis.set_major_formatter(ticker.PercentFormatter(1))
            plt.xlabel("")
            max_value = max(filtered_sheet.max())  # get the max value in this dataframe
            plt.ylim(0, max_value + 0.03)  # set y axis limit 1 - adjust frame
            plt.legend(bbox_to_anchor=(1.2, 1), loc='upper right')  # legend outside the graph
            # attach values
            for p in graph_sue_type.patches:
                if p.get_height() * 100 < 1:
                    h = "{:.1%}".format(p.get_height())
                    if h == "1.0%": h = "1%"
                else:
                    h = "{:.0%}".format(p.get_height())
                graph_sue_type.annotate(str(h), (p.get_x() + 0.025, p.get_height() + 0.005), size=8, rotation=45)
            plt.savefig(self.output_directory + '/Graphs/' + 'גרף_כמות_ילדים' + '.png',
                        bbox_inches='tight')  # save to folder as .png
            plt.clf()  # clear

        """Family situation"""
        if "מצב משפחתי" in self.sheet_pd:
            filtered_sheet = self.query_table_percents.copy()  # go over the df we made and brake it to tables
            filtered_sheet = filtered_sheet[
                (filtered_sheet['אלמנט השוואה'] == 'רווק/ה') | (
                        filtered_sheet['אלמנט השוואה'] == 'נשוי/ה') | (
                        filtered_sheet['אלמנט השוואה'] == 'גרוש/ה') | (
                        filtered_sheet['אלמנט השוואה'] == 'אלמנ/ה')]
            filtered_sheet.set_index('אלמנט השוואה', inplace=True)
            query_sum_arr_for_graphs.append(filtered_sheet)  # for result
            # legend names
            category_names_t = ["רווק/ה", " נשוי/ה", "גרוש/ה", "אלמנ/ה"]
            category_names = []
            for i in category_names_t: category_names.append(i[::-1])
            # set results
            results = {}
            for group in filtered_sheet.columns:
                if group != "פוליגם":
                    tmp = filtered_sheet.filter(items=[group])
                    y = tmp[group].to_numpy()  # get values
                    # add group as key with values y
                    results[group] = y
            # set labels - focus groups
            labels = []
            for word in list(filtered_sheet.columns): labels.append(word[::-1])  # Hebrew translation
            # build graph:
            data = np.array(list(results.values()))
            data_cum = data.cumsum(axis=1)
            category_colors = plt.get_cmap('Accent')(np.linspace(0.15, 0.85, data.shape[1]))
            fig, ax = plt.subplots(figsize=(9.2, 5))
            ax.invert_yaxis()
            ax.xaxis.set_visible(False)
            ax.set_xlim(0, np.sum(data, axis=1).max())
            for i, (colname, color) in enumerate(zip(category_names, category_colors)):
                widths = data[:, i]
                starts = data_cum[:, i] - widths
                ax.barh(labels, widths, left=starts, height=0.5,
                        label=colname, color=color)
                xcenters = starts + widths / 2
                text_color = 'black'
                for y, (x, c) in enumerate(zip(xcenters, widths)):
                    ax.text(x, y, str("{:.0%}".format(c)), ha='center', va='center',
                            color=text_color, fontsize=10, weight="bold")
            ax.legend(ncol=len(category_names), bbox_to_anchor=(0, 1),
                      loc='lower left', fontsize='small')
            plt.title("מצב משפחתי"[::-1], fontweight='bold')
            plt.savefig(self.output_directory + '/Graphs/' + 'גרף_מצב_משפחתי' + '.png',
                        bbox_inches='tight')  # save to folder as .png
            plt.clf()  # clear

        """Education"""
        if "רמת השכלה" in self.sheet_pd:
            filtered_sheet = self.query_table_percents.copy()  # go over the df we made and brake it to tables
            filtered_sheet = filtered_sheet[
                (filtered_sheet['אלמנט השוואה'] == 'תעודת מקצוע') | (
                        filtered_sheet['אלמנט השוואה'] == 'תואר שלישי') | (
                        filtered_sheet['אלמנט השוואה'] == 'תואר שני') | (
                        filtered_sheet['אלמנט השוואה'] == 'תואר ראשון') | (
                        filtered_sheet['אלמנט השוואה'] == 'יסודי/תיכון') | (
                        filtered_sheet['אלמנט השוואה'] == 'תעודת בגרות') | (
                        filtered_sheet['אלמנט השוואה'] == 'ללא השכלה')]
            filtered_sheet.set_index('אלמנט השוואה', inplace=True)
            query_sum_arr_for_graphs.append(filtered_sheet)  # for result
            colors = ['#ff9999', '#66b3ff', '#99ff99', '#ffcc99', '#E3C800', '#87794E', '#647687']
            i = 0  # create a different pie chart
            for group in filtered_sheet.columns:
                tmp = filtered_sheet.filter(items=[group])
                # delete if 0% from graph only
                tmp.drop(tmp.loc[tmp[group] < 0.01].index, inplace=True)
                y = tmp[group].to_numpy()  # get values
                # graph creation
                mylabels = list(tmp.index.values)  # crate labels
                mylabels_bckword = []
                for word in mylabels: mylabels_bckword.append(word[::-1])  # Hebrew translation
                plt.pie(y, labels=mylabels_bckword, colors=colors, shadow=True, autopct='%1.0f%%',
                        normalize=False)  # create the pie
                plt.title(("מצב השכלה: " + group)[::-1], fontweight='bold')
                plt.savefig(self.output_directory + '/Graphs/' + 'גרף_השכלה_' + str(group) + '.png',
                            bbox_inches='tight')  # save to folder as .png
                i += 1
                plt.clf()

        """Job AND Fields distribution"""
        for tmp_dict in [self.jobs_dic, self.fields_jobs_dic]:
            create_graph_jobs(self, tmp_dict, color_palette)  # w/ problematic values
            create_graph_jobs(self, tmp_dict, color_palette, True)  # w/o problematic values

        plt.close("all")
        return query_sum_arr_for_graphs

    def create_pptx(self):
        print("Creating Power point presentation")
        # name of pptx
        user_title_name = "השוואת "
        for name in self.filter_instructions_array:
            user_title_name = user_title_name + name[0]
            if name[0] != self.filter_instructions_array[-1][0]:
                user_title_name = user_title_name + ", "
        print("User name given for pptx: " + str(user_title_name))
        # open a new pptx
        try:
            prs = Presentation()  # open a presentation
        except:
            prs = Presentation('default_template.pptx')
        prs.slide_width = Inches(16)  # set slides sizes
        prs.slide_height = Inches(9)  # set slide sizes

        """Title slide"""
        if 1:
            # background
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # adding a slide + choosing a slide layout blank
            left = top = Inches(0)  # pic position
            img_path = "src_files/bck_first_slide.png"  # name of pic
            pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width,
                                           height=prs.slide_height)  # set background
            slide.shapes._spTree.remove(pic._element)  # This moves it to the background
            slide.shapes._spTree.insert(2, pic._element)  # This moves it to the background
            # title
            left = Inches(11.5)
            top = Inches(2)
            width = Inches(4)
            height = Inches(2)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "נתוני שירות התעסוקה"
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(64)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Ariel'
            # sub-title
            left = Inches(11.5)
            top = Inches(4)
            width = Inches(4)
            height = Inches(2)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            for group in self.filter_instructions_array:
                try:
                    text = text + group[0] + "\n"
                except UnboundLocalError:
                    text = group[0] + "\n"
            p.text = text
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Ariel'
            # date of today
            left = Inches(11.8)
            top = Inches(7.95)
            width = Inches(4)
            height = Inches(2)
            txBox = slide.shapes.add_textbox(left, top, width, height)  # right down corner
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = str(datetime.today().strftime('%d/%m/%Y'))
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(35)
            p.font.bold = True
            p.font.name = 'Ariel'

        """General details slide"""
        if 1:
            slide = new_body_slide(prs, "נתונים כלליים")

            # details - build the string
            temp = self.query_table_numbers.copy()  # go over the df we made and brake it to tables
            temp = temp[(temp['אלמנט השוואה'] == 'סכום כלל דורשי עבודה')].T
            temp_np = temp.to_numpy()
            i = 0
            details = ""
            for num in temp_np[1::]:
                num_insert = "{:,}".format(int(num[0]))
                name = self.filter_instructions_array[i][0]
                details = details + num_insert + "סך הכל דורשי עבודה ב" + name + "  " + "\n\n"
                i += 1
            # assign details
            left = Inches(3.9)
            top = Inches(2)
            width = Inches(12)
            height = Inches(4)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = details
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(35)
            p.font.bold = True

        """Sue type slide"""
        if "סוג תביעה נוכחי" in self.sheet_pd:
            slide = new_body_slide(prs, "התפלגות סוג תביעה")

            # add graphs
            pos = 0
            for group_set in self.filter_instructions_array[::-1]:
                name = group_set[0]
                img_path = self.output_directory + '/Graphs/' + 'גרף_סוג_תביעה_' + name + '.png'
                width = Inches(16 / len(self.filter_instructions_array))
                if len(self.filter_instructions_array) == 1:
                    width = Inches(5.5)
                    pos=5.25
                img = slide.shapes.add_picture(img_path, left=Inches(pos), top=Inches(1.5), width=width)
                pos = pos + 16 / len(self.filter_instructions_array)  # 16 is slide width

            # add a blank paragraph
            new_blank_paragraph(slide)

            # add the table
            img_path = self.output_directory + "/Tables/" + 'סוג תביעה' + '.png'
            img = slide.shapes.add_picture(img_path, left=Inches(0.2), top=Inches(7), width=Inches(5))

        """reason os registration slide"""
        if "סיבת רישום" in self.sheet_pd:
            slide = new_body_slide(prs, "התפלגות סיבת הרישום של דורשי עבודה לשירות")

            # add graph
            img_path = self.output_directory + "/Graphs/" + 'גרף_סיבת_רישום' + '.png'
            img = slide.shapes.add_picture(img_path, left=Inches(2), top=Inches(1.2), width=Inches(12))

            # add a blank paragraph
            new_blank_paragraph(slide)

            # add the table
            img_path = self.output_directory + "/Tables/" + 'סיבת רישום' + '.png'
            img = slide.shapes.add_picture(img_path, left=Inches(0.2), top=Inches(6.8), width=Inches(5))

        """Gender slide"""
        if "מגדר" in self.sheet_pd:
            slide = new_body_slide(prs, "התפלגות מגדרית")

            # add graph
            img_path = self.output_directory + '/Graphs/' + 'גרף_מגדר' + '.png'
            img = slide.shapes.add_picture(img_path, left=Inches(3), top=Inches(1.5), height=Inches(5.5))

            # add a blank paragraph
            new_blank_paragraph(slide)

            # add the table
            img_path = self.output_directory + "/Tables/" + 'מגדר' + '.png'
            img = slide.shapes.add_picture(img_path, left=Inches(0.2), top=Inches(7.4), width=Inches(5))

        """Age slide"""
        if "גיל" in self.sheet_pd:
            slide = new_body_slide(prs, "התפלגות הגילאים של דורשי עבודה")

            # add graph
            img_path = self.output_directory + '/Graphs/' + 'גרף_גילאים' + '.png'
            img = slide.shapes.add_picture(img_path, left=Inches(5.7), top=Inches(1.6), width=Inches(10))

            # add a blank paragraph
            new_blank_paragraph(slide)

            # add the table
            img_path = self.output_directory + "/Tables/" + 'גיל' + '.png'
            img = slide.shapes.add_picture(img_path, left=Inches(0.2), top=Inches(1.7), width=Inches(5))

        """Education slide"""
        if "רמת השכלה" in self.sheet_pd:
            slide = new_body_slide(prs, "התפלגות רמות ההשכלה של דורשי עבודה")

            # add graphs
            pos = 0  # position from left- alignment
            for group_set in self.filter_instructions_array[::-1]:
                name = group_set[0]
                img_path = self.output_directory + '/Graphs/' + 'גרף_השכלה_' + name + '.png'
                width = Inches(16 / len(self.filter_instructions_array))
                if len(self.filter_instructions_array) == 1:
                    width = Inches(5.5)
                    pos = 5.25
                img = slide.shapes.add_picture(img_path, left=Inches(pos), top=Inches(1.5), width=width)
                pos = pos + 16 / len(self.filter_instructions_array)

        # add a blank paragraph
        new_blank_paragraph(slide)

        # add the table
        img_path = self.output_directory + "/Tables/" + 'רמת השכלה' + '.png'
        img = slide.shapes.add_picture(img_path, left=Inches(0.2), top=Inches(6.3), width=Inches(4))

        """Family status slide"""
        if "מצב משפחתי" in self.sheet_pd:
            slide = new_body_slide(prs, "התפלגות מצב משפחתי")

            # add graph
            img_path = self.output_directory + '/Graphs/' + 'גרף_מצב_משפחתי' + '.png'
            img = slide.shapes.add_picture(img_path, left=Inches(3), top=Inches(1.5), height=Inches(5))

            # add a blank paragraph
            new_blank_paragraph(slide)

            # add the table
            img_path = self.output_directory + "/Tables/" + 'מצב משפחתי' + '.png'
            img = slide.shapes.add_picture(img_path, left=Inches(0.2), top=Inches(7), width=Inches(5))

        """children quantity slide"""
        if "ילדים עד גיל 18" in self.sheet_pd:
            slide = new_body_slide(prs, '18' + "התפלגות כמות ילדים מתחת לגיל ")

            # add graph
            img_path = self.output_directory + '/Graphs/' + 'גרף_כמות_ילדים' + '.png'
            img = slide.shapes.add_picture(img_path, left=Inches(3.5), top=Inches(1.5), height=Inches(5))
            # add a blank paragraph
            new_blank_paragraph(slide)

            # add the table
            img_path = self.output_directory + "/Tables/" + 'כמות ילדים' + '.png'
            img = slide.shapes.add_picture(img_path, left=Inches(0.2), top=Inches(6.7), width=Inches(4))

        """Job distribution - 3 slides"""
        if "מקצועות רלוונטיים" in self.sheet_pd:
            # slide 1
            if 1:
                text = "מקצועות שכיחים לפי דירוג " + str(self.filter_instructions_array[-1][0])
                slide = new_body_slide(prs, text)

                # add table
                img_path = self.output_directory + '/Tables/' + 'מקצועות שכיחים' + '.png'
                img = slide.shapes.add_picture(img_path, left=Inches(4), top=Inches(1.5), width=Inches(8))

            # slide 2
            if 1:
                slide = new_body_slide(prs, "התפלגות משלחי היד - מבט כללי")

                # add graph
                img_path = self.output_directory + '/Graphs/' + 'גרף_מקצועות' + '.png'
                img = slide.shapes.add_picture(img_path, left=Inches(2), top=Inches(1.3), height=Inches(7.5))

            # slide 3
            if 1:
                slide = new_body_slide(prs, "התפלגות משלחי היד - מבט כללי")

                # add graph
                img_path = self.output_directory + '/Graphs/' + 'גרף_מקצועות_ללא_כללי' + '.png'
                img = slide.shapes.add_picture(img_path, left=Inches(2), top=Inches(1.3), height=Inches(7.5))

        """Fields Job distribution - 3 slides"""
        if "ענפי מקצועות רלוונטיים" in self.sheet_pd:
            # slide 1
            if 1:
                text = "ענפי מקצועות שכיחים לפי דירוג " + str(self.filter_instructions_array[-1][0])
                slide = new_body_slide(prs, text)

                # add table
                img_path = self.output_directory + '/Tables/' + 'ענפי מקצועות שכיחים' + '.png'
                img = slide.shapes.add_picture(img_path, left=Inches(4), top=Inches(1.5), width=Inches(8))

            # slide 2
            if 1:
                slide = new_body_slide(prs, "התפלגות ענפי משלחי היד - מבט כללי")

                # add graph
                img_path = self.output_directory + '/Graphs/' + 'גרף_ענפים' + '.png'
                img = slide.shapes.add_picture(img_path, left=Inches(2), top=Inches(1.3), height=Inches(7.5))

            # slide 3
            if 1:
                slide = new_body_slide(prs, "התפלגות ענפי משלחי היד - מבט כללי")

                # add graph
                img_path = self.output_directory + '/Graphs/' + 'גרף_ענפים_ללא_כללי' + '.png'
                img = slide.shapes.add_picture(img_path, left=Inches(2), top=Inches(1.3), height=Inches(7.5))

        """ Finish slide """
        if 1:
            slide = new_body_slide(prs, "שורות תחתונות")
            new_blank_paragraph(slide, top=Inches(2))

        # save to file
        prs.save(self.output_directory + "/" + user_title_name + ".pptx")  # saving file

    def create_excel_sum_ups(self):
        print("Exporting data to EXCEL")

        with pd.ExcelWriter(r'' + self.output_directory + "\Output_Report.xlsx") as writer:

            # SHEET 1: Result numbers
            self.query_table_numbers.to_excel(writer, sheet_name='QueryResults-count', index=False)
            workbook = writer.book

            # SHEET 2: Result precents
            self.query_table_percents.to_excel(writer, sheet_name='QueryResults-percent', index=False)

            # Design sheets #1 and #2
            note = str(get_max_column_excel(self))  # get the last letter to put background on
            for name in ['QueryResults-count', 'QueryResults-percent']:
                worksheet = writer.sheets[name]  # get to the sheet
                worksheet.right_to_left()
                header_fmt_r = workbook.add_format({'align': 'right', 'bold': True, 'font_color': '#2B337A'})
                worksheet.set_column('A:A', 25, header_fmt_r)
                header_fmt = workbook.add_format({'align': 'center', 'bold': True, 'font_color': '#2B337A'})
                worksheet.set_row(2, 15, header_fmt)  # make headers of groups in bold and color blue
                worksheet.set_row(6, 15, header_fmt)
                worksheet.set_row(11, 15, header_fmt)
                worksheet.set_row(17, 15, header_fmt)
                worksheet.set_row(30, 15, header_fmt)
                worksheet.set_row(37, 15, header_fmt)
                worksheet.set_row(42, 15, header_fmt)
                worksheet.set_row(51, 15, header_fmt)
                worksheet.set_row(52, 15, header_fmt)
                data_fmt = workbook.add_format({'align': 'center'})
                worksheet.set_column('B:' + note, 15, data_fmt)
                worksheet.insert_image(str(chr(ord(note) + 1)) + '1', 'src_files/icon.png')
                # background
                format_bck_orange = workbook.add_format({'bg_color': '#FFBF48', 'border': 1})
                format_bck_dark_blue_font_white = workbook.add_format({'bg_color': '#002060', 'font_color': '#FFFFFF'})
                # headers
                worksheet.conditional_format('A3',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_dark_blue_font_white})
                worksheet.conditional_format('A7',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_dark_blue_font_white})
                worksheet.conditional_format('A12',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_dark_blue_font_white})
                worksheet.conditional_format('A18',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_dark_blue_font_white})
                worksheet.conditional_format('A31',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_dark_blue_font_white})
                worksheet.conditional_format('A38',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_dark_blue_font_white})
                worksheet.conditional_format('A43',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_dark_blue_font_white})
                worksheet.conditional_format('A3:' + note + '5',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_orange})
                worksheet.conditional_format('A7:' + note + '10',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_orange})
                worksheet.conditional_format('A12:' + note + '16',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_orange})
                worksheet.conditional_format('A18:' + note + '29',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_orange})
                worksheet.conditional_format('A31:' + note + '36',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_orange})
                worksheet.conditional_format('A38:' + note + '41',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_orange})
                worksheet.conditional_format('A43:' + note + '50',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_orange})
                worksheet.conditional_format('A52:' + note + '53',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0,
                                              'format': format_bck_orange})
                if name == 'QueryResults-percent':
                    percent_fmt = workbook.add_format({'num_format': '0.0%', 'align': 'center'})
                    worksheet.set_column('B:' + note, 15, percent_fmt)  # Quota percent columns

            # SHEET 3: Result Jobs
            if "מקצועות רלוונטיים" in self.sheet_pd:
                self.jobs_dic.sort_values(by=self.filter_instructions_array[-1][0], ascending=False).to_excel(writer,
                                                                                                              sheet_name='Jobs distribution',
                                                                                                              index=True)
                worksheet = writer.sheets['Jobs distribution']  # get to the sheet
                worksheet.right_to_left()
                worksheet.set_column('A:A', 35)
                worksheet.set_column('B:' + get_max_column_excel(self), 15)
                format_border = workbook.add_format({'border': 1, 'bg_color': '#73b7ff'})
                format_green = workbook.add_format({'border': 1, 'bg_color': '#bcffad'})
                l = str(len(self.jobs_dic.index) + 1)
                print("Number of Jobs found: " + str(len(self.jobs_dic.index)))
                worksheet.conditional_format('A1:' + note + '1',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0, 'format': format_border})
                worksheet.conditional_format('A1:' + note + l,
                                             {'type': 'cell', 'criteria': '>=', 'value': 1, 'format': format_green})
                worksheet.insert_image(str(chr(ord(get_max_column_excel(self)) + 1)) + '1', 'src_files/icon.png')

            # SHEET 3: Result Jobs
            if "ענפי מקצועות רלוונטיים" in self.sheet_pd:
                self.fields_jobs_dic.sort_values(by=self.filter_instructions_array[-1][0], ascending=False).to_excel(
                    writer,
                    sheet_name='Fields distribution', index=True)
                worksheet = writer.sheets['Fields distribution']  # get to the sheet
                worksheet.right_to_left()
                worksheet.set_column('A:A', 35)
                worksheet.set_column('B:' + get_max_column_excel(self), 15)
                format_border = workbook.add_format({'border': 1, 'bg_color': '#73b7ff'})
                format_green = workbook.add_format({'border': 1, 'bg_color': '#bcffad'})
                l = str(len(self.jobs_dic.index) + 1)
                print("Number of Jobs found: " + str(len(self.jobs_dic.index)))
                worksheet.conditional_format('A1:' + note + '1',
                                             {'type': 'cell', 'criteria': '>=', 'value': 0, 'format': format_border})
                worksheet.conditional_format('A1:' + note + l,
                                             {'type': 'cell', 'criteria': '>=', 'value': 1, 'format': format_green})
                worksheet.insert_image(str(chr(ord(get_max_column_excel(self)) + 1)) + '1', 'src_files/icon.png')

            print("Saving report sum up to EXCEL... takes a few moments")

        pass


class StandardAnalysis_byID(StandardAnalysis):

    def __init__(self, id_sheet=None):
        StandardAnalysis.__init__(self, None, None, None, None, None, None)
        self.id_sheet = id_sheet

    def get_ID_filtered_sheet(self, id_excel):
        print("Getting the dataframe for the list of IDs")
        # first step - get ID list
        sheet_IDs = get_sheet_pd(id_excel)
        tmp = list(sheet_IDs)
        head_name = ""
        if 'ת"ז' in tmp:
            head_name = 'ת"ז'
        elif "תעודת זהות" in tmp:
            head_name = "תעודת זהות"
        elif "ID" in tmp:
            head_name = "ID"
        elif "תז" in tmp:
            head_name = "תז"
        elif "מספר זהות" in tmp:
            head_name = "מספר זהות"
        elif "מספר מזהה" in tmp:
            head_name = "מספר מזהה"
        else:
            print("NO ID field identifier was found.\nSorry but the program will close now.")
            alert_popup("Error", "לא נמצא שדה תעודת זהות (אנא שנה ידנית את השדה בקובץ אצלך לשם 'מספר זהות'")
            exit(1)
        print("Found a field for ID: " + str(head_name))
        list_IDs = sheet_IDs[head_name].unique()
        print("There are " + str(len(list_IDs)) + " IDs in the secondary input file.")
        print("There are " + str(self.sheet_pd.shape[0]) + " inputs in the original input file.")
        # second step - filter our dataframe through the IDs
        try:
            self.id_sheet = self.sheet_pd[self.sheet_pd[head_name].isin(list_IDs)]
        except KeyError:
            print("NO IDs in the input file.\nSorry but the program will close now.")
            alert_popup("Error", "לא נמצאו תעודות זהות בקובץ הנתונים שהזנת")
            exit(1)

    def get_dictionary(self, param):
        """
            This function tops the Itzik function and after splitting - summing to a dictionary
            :param param: either jobs or fields (hebrew string)
            :return: a python dataframe that came from  dictionary that sums the param through all of the dataframe
        """
        if param in self.sheet_pd:
            print("")  # create separation
            sheet = get_splitted_sheet(self.sheet_pd, param)
            sheet_ids = get_splitted_sheet(self.id_sheet, param)
            # create a comparison of job distributions
            heads = []
            l = 0
            output_dic = {}
            for col in self.filter_instructions_array:
                heads.append(col[0])  # create array heads for focus groups
                temp_filtered_df = sheet_pd_filter(sheet, col[1])  # filter dataframe
                # keep only columns with jobs
                pool_jobs = temp_filtered_df[
                    [str(param) + " " + str(1), str(param) + " " + str(2), str(param) + " " + str(3),
                     str(param) + " " + str(4)]]
                # create a dictionary {key = str of job : value = arr of int length of focus groups }
                output_dic = update_dic_values(output_dic, pool_jobs, l, len(self.filter_instructions_array) + 1)
                l = l + 1  # add to iterator
            # add the extra column for the IDs
            heads.append("רשימת ת''ז")  # create array heads for focus groups
            temp_filtered_df = sheet_ids  # filter dataframe
            # keep only columns with jobs
            pool_jobs = temp_filtered_df[
                [str(param) + " " + str(1), str(param) + " " + str(2), str(param) + " " + str(3),
                 str(param) + " " + str(4)]]
            # create a dictionary {key = str of job : value = arr of int length of focus groups }
            output_dic = update_dic_values(output_dic, pool_jobs, l, len(self.filter_instructions_array) + 1)
            # return a data frame
            if param == "מקצועות רלוונטיים":
                self.jobs_dic = pd.DataFrame.from_dict(output_dic, orient='index', columns=heads)
            elif param == "ענפי מקצועות רלוונטיים":
                self.fields_jobs_dic = pd.DataFrame.from_dict(output_dic, orient='index', columns=heads)
        else:
            print("No " + str(param) + " in the input excel sheet. Not activating the Itzik function.")
            return None

    def set_query_tables(self):
        print("Counting query data, setting query tables.")
        # prepare to count - add rows - IES format!
        rows_names = [
            None, "מגדר",
            "נקבה", "זכר",
            None, "סוג תביעה",
            "אבטלה", "הבטחת הכנסה", "אינו תובע",
            None, "סיבת רישום",
            "פיטורין", "חל''ת", "אינו עובד ומחפש עבודה", "התפטרות",
            None, "גיל",
            "15-24", "25-29", "30-34", "35-39", "40-44", "45-49", "50-54", "55-59", "60-64", "65-69", "70 +",
            None, "מספר ילדים עד גיל 18",
            "0", "1-2", "3-5", "6-8", "יותר מ-8",
            None, "מצב משפחתי",
            "רווק/ה", "נשוי/ה", "גרוש/ה",
            None, "מצב השכלה",
            "ללא השכלה", "תעודת בגרות", "יסודי/תיכון", "תואר ראשון", "תואר שני", "תואר שלישי", "תעודת מקצוע",
            None, None,
            "סכום כלל דורשי עבודה"
        ]
        output_num = pd.DataFrame({"אלמנט השוואה": rows_names}, columns=["אלמנט השוואה"])
        output_per = pd.DataFrame({"אלמנט השוואה": rows_names}, columns=["אלמנט השוואה"])
        # start counting - go over database {no. of columns} times
        for col in self.filter_instructions_array:
            name_col = col[0]
            temp_filtered_df = sheet_pd_filter(self.sheet_pd, col[1])  # filter dataframe
            # add output to
            temp_output = query_counter_helper(temp_filtered_df, name_col)  # counter helper
            output_num[name_col] = temp_output[0]
            output_per[name_col] = temp_output[1]
        # additional column for IDs
        name_col = "רשימת ת''ז"
        temp_output = query_counter_helper(self.id_sheet, name_col)  # counter helper
        output_num[name_col] = temp_output[0]
        output_per[name_col] = temp_output[1]
        # finish
        self.query_table_numbers = output_num
        self.query_table_percents = output_per


class AutoAnalysis:

    # instance attributes
    def __init__(self, sheet_pd=None, filter_instructions_array=None, output_directory=None):
        self.sheet_pd = sheet_pd
        self.filter_instructions_array = filter_instructions_array
        self.output_directory = output_directory
        self.matrix_result_arr = []  # node - [name of group , the matrix as array ]

    def matrix_creator(self):
        print("Creating the 2-D matrix.")
        matrix_result_arr = []  # [name of group , the matrix as array ]
        # fill values in the matrix for each group
        for group in self.filter_instructions_array:
            group_name = group[0]
            filtered_df = sheet_pd_filter(self.sheet_pd, group[1])
            matrix_group = []
            # create the frame of the matrix
            headers = list(self.sheet_pd)  # big headlines
            indexes_matrix = []
            indexes_values = []
            indexes_values_specials = []
            for header in headers:
                values_header = filtered_df[header].unique()
                # take only those who have countable values
                if len(values_header) <= 15:
                    for value in values_header:
                        if str(value) != 'nan':
                            indexes_matrix.append(' - '.join([str(header), str(value)]))
                            indexes_values.append([header, value])
                # exceptions for count
                else:
                    if len(values_header) < 50:
                        try:
                            # check if numbers
                            for value in list(filtered_df[header].value_counts().iloc[0:15].index):
                                if int(value) != 'nan':
                                    indexes_matrix.append(' - '.join([str(header), str(value)]))
                                    indexes_values.append([header, value])
                        except ValueError:
                            # if < 50 && not number then take first 15 values
                            for value in list(filtered_df[header].value_counts().iloc[0:15].index):
                                if str(value) != 'nan':
                                    indexes_matrix.append(' - '.join([str(header), str(value)]))
                                    indexes_values.append([header, value])
            i = j = 0
            # make the matrix frame
            for pair1 in tqdm(indexes_values):  # row
                matrix_group.append([])
                for pair2 in indexes_values:  # col
                    count = \
                        filtered_df[(filtered_df[pair1[0]] == pair1[1]) & (filtered_df[pair2[0]] == pair2[1])].shape[0]
                    matrix_group[i].append(count)
                    j += 1
                i += 1

            # append to the array of focus groups
            matrix_result_arr.append(
                [group_name, pd.DataFrame(matrix_group, columns=indexes_matrix, index=indexes_matrix)])
        self.matrix_result_arr = matrix_result_arr
        # open a new excel
        with pd.ExcelWriter(r'' + self.output_directory + "\Output_Crossing_Data.xlsx") as writer:
            # each sheet is a different matrix
            for i in range(len(matrix_result_arr)):
                matrix_result_arr[i][1].to_excel(writer, sheet_name=matrix_result_arr[i][0], index=True)
                workbook = writer.book
                worksheet = writer.sheets[matrix_result_arr[i][0]]  # get to the sheet
                worksheet.right_to_left()
                header_fmt_r = workbook.add_format({'align': 'right', 'bold': True, 'font_color': '#2B337A'})
                worksheet.set_column('A:A', 30, header_fmt_r)
                # background
                format_bck_green = workbook.add_format({'bg_color': '#CCFF99', 'border': 1})
                format_bck_orange = workbook.add_format({'bg_color': '#FFCC99', 'border': 1})
                format_bck_yellow = workbook.add_format({'bg_color': '#FFFF99', 'border': 1})
                format_bck_gray = workbook.add_format({'bg_color': '#E0E0E0', 'border': 1})
                worksheet.conditional_format('B2:DA100', {'type': 'cell', 'criteria': '>', 'value': 200,
                                                          'format': format_bck_green})
                worksheet.conditional_format('B2:DA100',
                                             {'type': 'cell', 'criteria': 'between', 'minimum': 30, 'maximum': 199,
                                              'format': format_bck_orange})
                worksheet.conditional_format('B2:DA100',
                                             {'type': 'cell', 'criteria': 'between', 'minimum': 1, 'maximum': 29,
                                              'format': format_bck_yellow})
                worksheet.conditional_format('B2:DA100', {'type': 'bottom', 'value': "1", 'format': format_bck_gray})

    def query_creator(self):
        print("Creating the query tables.")
        filtered_df = sheet_pd_filter(self.sheet_pd, self.filter_instructions_array[0][1])
        # create the lists of the summs up
        headers = list(self.sheet_pd)  # big headlines
        indexes_query = []
        indexes_values = []
        for header in headers:
            values_header = filtered_df[header].unique()
            # take only those who have countable values
            if len(values_header) <= 20:
                for value in values_header:
                    if str(value) != 'nan':
                        indexes_query.append('-'.join([str(header), str(value)]))
                        indexes_values.append([header, value])
            # exceptions for count
            else:
                if len(values_header) < 80:
                    try:
                        # check if numbers
                        for value in values_header:
                            if int(value) != 'nan':
                                indexes_query.append(' - '.join([str(header), str(value)]))
                                indexes_values.append([header, value])
                    except ValueError:
                        # if < 50 && not number then take first 15 values
                        for value in list(filtered_df[header].value_counts().iloc[0:15].index):
                            if str(value) != 'nan':
                                indexes_query.append(' - '.join([str(header), str(value)]))
                                indexes_values.append([header, value])
        # prepare the filtered df of the focus groups
        df_grups = []
        for group in self.filter_instructions_array:
            df_grups.append(sheet_pd_filter(self.sheet_pd, group[1]))
        # make the query frame
        res = []
        i = 0
        for pair in tqdm(indexes_values):  # row
            res.append([])
            for df in df_grups:  # col
                count = df[df[pair[0]] == pair[1]].shape[0]
                res[i].append(count)
            i += 1
        # finish
        res = pd.DataFrame(res, columns=[x[0] for x in self.filter_instructions_array], index=indexes_query)
        # open a new excel
        with pd.ExcelWriter(r'' + self.output_directory + "\Output_Query_Data.xlsx") as writer:
            res.to_excel(writer, sheet_name="query_result", index=True)
            workbook = writer.book
            worksheet = writer.sheets["query_result"]  # get to the sheet
            worksheet.right_to_left()
            header_fmt_r = workbook.add_format({'align': 'right', 'bold': True, 'font_color': '#2B337A'})
            worksheet.set_column('A:A', 30, header_fmt_r)
            worksheet.set_column('B:' + str(get_max_column_excel(self)), 20)
            # background
            format_bck_green = workbook.add_format({'bg_color': '#CCFF99', 'border': 1})
            format_bck_orange = workbook.add_format({'bg_color': '#FFCC99', 'border': 1})
            format_bck_yellow = workbook.add_format({'bg_color': '#FFFF99', 'border': 1})
            format_bck_gray = workbook.add_format({'bg_color': '#E0E0E0', 'border': 1})
            worksheet.conditional_format('B2:DA1000', {'type': 'cell', 'criteria': '>', 'value': 200,
                                                       'format': format_bck_green})
            worksheet.conditional_format('B2:DA1000',
                                         {'type': 'cell', 'criteria': 'between', 'minimum': 30, 'maximum': 199,
                                          'format': format_bck_orange})
            worksheet.conditional_format('B2:DA1000',
                                         {'type': 'cell', 'criteria': 'between', 'minimum': 1, 'maximum': 29,
                                          'format': format_bck_yellow})
            worksheet.conditional_format('B2:DA1000', {'type': 'bottom', 'value': "1", 'format': format_bck_gray})
